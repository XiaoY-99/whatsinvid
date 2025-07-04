from pptx import Presentation
from pptx.util import Inches
import os

def generate_slides(summary: str, output_path: str) -> str:
    """
    Generate a simple PowerPoint file from a translated summary, each with a logo.

    Args:
        summary (str): Translated summary (multiline).
        output_path (str): Path to save the PPTX file.

    Returns:
        str: Path to the generated slide file.
    """
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    # Path to logo image (ensure it's in the same dir as this script)
    logo_path = os.path.join(os.path.dirname(__file__), "logo.png")
    add_logo = os.path.exists(logo_path)

    for idx, bullet in enumerate(summary.split('\n')):
        if bullet.strip():
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]

            title.text = f"Point {idx + 1}"
            content.text = bullet.strip()

            # Add logo to slide (e.g., bottom-right corner)
            if add_logo:
                slide.shapes.add_picture(
                    logo_path,
                    left=prs.slide_width - Inches(1.8),     # near the right edge
                    top=prs.slide_height - Inches(1.2),     # near the bottom
                    width=Inches(1.5)                       # adjust size as needed
                )

    prs.save(output_path)
    return output_path
